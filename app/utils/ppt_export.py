"""
Module d'export en PowerPoint pour AutoCours.
Utilise python-pptx pour générer un fichier .pptx à partir du contenu du cours.
"""

from pptx import Presentation

def export_to_ppt(course_content: dict) -> str:
    """
    Exporte le contenu du cours dans un fichier PowerPoint (.pptx).

    Args:
        course_content (dict): Dictionnaire contenant le contenu du cours.

    Returns:
        str: Chemin du fichier PowerPoint généré.
    """
    ppt_file = f"{course_content['title'].replace(' ', '_')}.pptx"
    prs = Presentation()

    # Création de la slide de titre
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = course_content["title"]
    slide.placeholders[1].text = "Généré automatiquement par AutoCours"

    # Création des slides pour chaque module
    for module in course_content["modules"]:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = module["title"]
        slide.placeholders[1].text = module["content"]

    prs.save(ppt_file)
    return ppt_file
